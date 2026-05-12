from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)   # 배경/제목
BLUE   = RGBColor(0x1A, 0x5F, 0xA8)   # 포인트
SKYB   = RGBColor(0x41, 0x9B, 0xD2)   # 보조
ORANGE = RGBColor(0xE8, 0x6A, 0x2B)   # 강조
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)   # 카드 배경
DGRAY  = RGBColor(0x44, 0x44, 0x55)   # 본문 텍스트
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃

# ── 헬퍼 ────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_bg(slide, color=NAVY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def slide_header(slide, title, subtitle=None):
    """상단 헤더 바 + 제목"""
    add_rect(slide, 0, 0, 13.33, 1.1, fill=NAVY)
    add_text(slide, title, 0.4, 0.15, 10, 0.8,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 10, 0.4,
                 size=11, color=SKYB)
    # 우측 로고 텍스트
    add_text(slide, "한국도로공사", 11.5, 0.25, 1.7, 0.5,
             size=11, color=RGBColor(0xA0,0xC0,0xE0), align=PP_ALIGN.RIGHT)

def kpi_card(slide, l, t, w, h, label, value, unit="", sub="", accent=BLUE):
    add_rect(slide, l, t, w, h, fill=WHITE,
             line=accent, line_w=Pt(2))
    add_rect(slide, l, t, w, 0.08, fill=accent)          # 상단 컬러 바
    add_text(slide, label, l+0.12, t+0.13, w-0.2, 0.35,
             size=11, color=DGRAY, bold=False)
    add_text(slide, value, l+0.12, t+0.45, w-0.2, 0.6,
             size=26, bold=True, color=accent)
    if unit:
        add_text(slide, unit, l+0.12+len(value)*0.18, t+0.72, w-0.3, 0.35,
                 size=11, color=DGRAY)
    if sub:
        add_text(slide, sub, l+0.12, t+h-0.45, w-0.2, 0.35,
                 size=10, color=RGBColor(0x80,0x80,0x95), italic=True)

def bar_h(slide, l, t, w, h, value, max_val, color=BLUE, bg=LGRAY):
    add_rect(slide, l, t, w, h, fill=bg)
    bar_w = w * value / max_val if max_val else 0
    add_rect(slide, l, t, bar_w, h, fill=color)

# ════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
# 배경 그라디언트 효과 (두 레이어)
add_rect(s1, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s1, 0, 0, 13.33, 3.5, fill=RGBColor(0x08,0x1C,0x3B))

# 상단 장식 바
add_rect(s1, 0, 3.35, 13.33, 0.08, fill=BLUE)
add_rect(s1, 0, 3.43, 13.33, 0.04, fill=ORANGE)

# 좌측 강조 선
add_rect(s1, 0.5, 1.2, 0.07, 1.8, fill=ORANGE)

# 제목
add_text(s1, "노선별 정체구간", 0.75, 1.1, 11, 1.0,
         size=44, bold=True, color=WHITE)
add_text(s1, "분석 보고서", 0.75, 2.0, 11, 0.9,
         size=44, bold=True, color=SKYB)

# 부제목
add_text(s1, "2026년 4월 고속도로 정체 현황 심층 분석",
         0.75, 3.0, 10, 0.5, size=16, color=RGBColor(0xA0,0xC0,0xE0))

# 하단 정보
add_rect(s1, 0, 5.8, 13.33, 1.7, fill=RGBColor(0x05,0x12,0x26))
add_text(s1, "한국도로공사  |  교통분석팀",
         0.6, 6.0, 8, 0.5, size=13, color=RGBColor(0x80,0xA0,0xC0))
add_text(s1, "2026. 05. 11",
         10.5, 6.0, 2.5, 0.5, size=13, color=RGBColor(0x80,0xA0,0xC0),
         align=PP_ALIGN.RIGHT)

# KPI 미리보기 카드 (3개)
cards = [("총 정체 건수","120건"),("평균 지속시간","60분"),("최대 정체길이","14.9km")]
for i,(lbl,val) in enumerate(cards):
    cx = 1.0 + i * 3.8
    add_rect(s1, cx, 4.5, 3.2, 1.1, fill=RGBColor(0x1A,0x3A,0x65))
    add_rect(s1, cx, 4.5, 3.2, 0.06, fill=ORANGE if i==2 else BLUE)
    add_text(s1, lbl, cx+0.15, 4.6, 2.9, 0.35, size=11,
             color=RGBColor(0xA0,0xC0,0xE0))
    add_text(s1, val, cx+0.15, 4.9, 2.9, 0.55, size=22, bold=True, color=WHITE)

# ════════════════════════════════════════════════════════
# 슬라이드 2 — 분석 개요
# ════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(s2, "분석 개요", "데이터 기간: 2026년 4월 1일 ~ 30일")
add_rect(s2, 0, 1.1, 13.33, 0.04, fill=BLUE)

# KPI 카드 4개
kpis = [
    ("총 정체 발생 건수", "120", "건", "4월 한 달간", BLUE),
    ("평균 지속시간",     "60.1","분", "최단 15분 / 최장 180분", ORANGE),
    ("평균 정체길이",     "3.7", "km","최대 14.9km", SKYB),
    ("분석 대상 노선",    "5",   "개", "경부·영동·서해안·중부·남해", GREEN),
]
for i,(lbl,val,unit,sub,col) in enumerate(kpis):
    cx = 0.35 + i * 3.22
    kpi_card(s2, cx, 1.3, 3.0, 1.6, lbl, val, unit, sub, col)

# 데이터 구성 설명
add_rect(s2, 0.35, 3.15, 12.6, 3.9, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_rect(s2, 0.35, 3.15, 12.6, 0.45, fill=BLUE)
add_text(s2, "데이터 항목 구성", 0.55, 3.18, 8, 0.4, size=13, bold=True, color=WHITE)

cols_info = [
    ("발생일시",   "정체 시작 시각 (연-월-일 시:분)"),
    ("노선",      "고속도로 노선명 (5개 노선)"),
    ("구간",      "IC/JC 기준 정체 구간명"),
    ("구간길이_km","해당 구간 전체 길이"),
    ("정체길이_km","실제 정체 발생 구간 길이"),
    ("지속시간_분", "정체 시작 ~ 해소까지 소요 시간"),
    ("원인",      "정체 유발 원인 (6가지 분류)"),
    ("해소시각",   "정체 해소 시각"),
]
for i,(col,desc) in enumerate(cols_info):
    row = i % 4
    c   = i // 4
    bx  = 0.55 + c * 6.2
    by  = 3.75 + row * 0.72
    add_rect(s2, bx, by, 1.35, 0.45, fill=RGBColor(0xE8,0xF0,0xFA))
    add_text(s2, col, bx+0.08, by+0.06, 1.2, 0.38, size=11, bold=True, color=BLUE)
    add_text(s2, desc, bx+1.45, by+0.06, 4.5, 0.38, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════
# 슬라이드 3 — 노선별 분석
# ════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(s3, "노선별 분석", "발생 건수 · 평균 지속시간 · 평균 정체길이")
add_rect(s3, 0, 1.1, 13.33, 0.04, fill=BLUE)

routes = [
    ("경부고속도로",  48, 53.0, 4.0,  14.9, BLUE),
    ("영동고속도로",  27, 62.9, 3.2,  12.9, SKYB),
    ("서해안고속도로",23, 69.5, 4.2,  12.3, GREEN),
    ("중부고속도로",  13, 53.8, 3.5,  14.7, ORANGE),
    ("남해고속도로",   9, 74.2, 2.7,   4.9, RGBColor(0x8E,0x44,0xAD)),
]

# 헤더
add_rect(s3, 0.35, 1.25, 12.6, 0.45, fill=NAVY)
for x,lbl in [(0.5,"노선"),(3.2,"발생 건수"),(5.5,"비율"),(7.2,"평균 지속(분)"),(9.5,"평균 정체(km)"),(11.5,"최대 정체(km)")]:
    add_text(s3, lbl, x, 1.3, 2.0, 0.38, size=11, bold=True, color=WHITE)

for i,(name,cnt,dur,avg_l,max_l,col) in enumerate(routes):
    by = 1.75 + i * 0.9
    bg = WHITE if i%2==0 else RGBColor(0xF7,0xF9,0xFC)
    add_rect(s3, 0.35, by, 12.6, 0.85, fill=bg)
    add_rect(s3, 0.35, by, 0.06, 0.85, fill=col)

    add_text(s3, name,            0.55, by+0.2, 2.5, 0.45, size=13, bold=True, color=DGRAY)
    add_text(s3, str(cnt),        3.2,  by+0.2, 1.5, 0.45, size=20, bold=True, color=col)
    add_text(s3, "건",            3.75, by+0.35, 0.5, 0.3, size=11, color=DGRAY)

    pct = round(cnt/120*100,1)
    add_text(s3, f"{pct}%",       5.5,  by+0.2, 1.5, 0.45, size=14, bold=True, color=DGRAY)

    # 건수 바
    bar_h(s3, 6.2, by+0.28, 2.6, 0.3, cnt, 48, color=col)
    add_text(s3, f"{dur}분",      9.5,  by+0.2, 1.8, 0.45, size=14, bold=True, color=DGRAY)
    add_text(s3, f"{avg_l}km",   11.2,  by+0.2, 1.5, 0.45, size=14, bold=True, color=DGRAY)

# 인사이트 박스
add_rect(s3, 0.35, 6.45, 12.6, 0.75, fill=RGBColor(0xE8,0xF0,0xFA), line=BLUE, line_w=Pt(1))
add_text(s3, "KEY INSIGHT",  0.55, 6.5, 1.8, 0.35, size=10, bold=True, color=BLUE)
add_text(s3,
    "경부고속도로가 전체의 40%를 차지하며 압도적 1위. 남해고속도로는 건수는 적지만 평균 지속시간 74분으로 최장 — 발생 시 장기화 위험.",
    2.2, 6.5, 10.5, 0.6, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════
# 슬라이드 4 — 원인별 분석
# ════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(s4, "원인별 분석", "정체 유발 원인 6가지 분류별 현황")
add_rect(s4, 0, 1.1, 13.33, 0.04, fill=BLUE)

causes = [
    ("공사",      25, 20.8, 54.9, 3.2,  BLUE),
    ("행사",      23, 19.2, 67.3, 4.9,  ORANGE),
    ("교통량 증가",23, 19.2, 49.3, 2.6,  SKYB),
    ("기상악화",   20, 16.7, 63.6, 3.6,  GREEN),
    ("사고",      18, 15.0, 59.9, 4.2,  RED),
    ("주말 정체",  11,  9.2, 73.2, 4.2,  RGBColor(0x8E,0x44,0xAD)),
]

# 왼쪽: 막대 차트형 카드
add_rect(s4, 0.35, 1.2, 7.5, 5.9, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_text(s4, "발생 건수 및 비율", 0.55, 1.25, 5, 0.4, size=13, bold=True, color=DGRAY)

for i,(name,cnt,pct,dur,avg_l,col) in enumerate(causes):
    by = 1.75 + i * 0.85
    add_text(s4, name, 0.55, by+0.1, 1.6, 0.4, size=12, bold=True, color=DGRAY)
    bar_h(s4, 2.25, by+0.1, 4.2, 0.38, cnt, 25, color=col)
    add_text(s4, f"{cnt}건 ({pct}%)", 6.55, by+0.1, 1.2, 0.38, size=11, color=col, bold=True)

# 오른쪽: 지속시간 비교 카드
add_rect(s4, 8.1, 1.2, 5.0, 2.8, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_text(s4, "원인별 평균 지속시간 (분)", 8.3, 1.25, 4.5, 0.4, size=13, bold=True, color=DGRAY)
for i,(name,cnt,pct,dur,avg_l,col) in enumerate(causes):
    by = 1.75 + i * 0.4
    add_text(s4, name, 8.3, by, 1.7, 0.35, size=10, color=DGRAY)
    bar_h(s4, 10.1, by+0.04, 2.5, 0.26, dur, 80, color=col)
    add_text(s4, f"{dur}분", 12.65, by, 0.6, 0.35, size=10, bold=True, color=col)

# 오른쪽: 정체길이 비교 카드
add_rect(s4, 8.1, 4.2, 5.0, 2.8, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_text(s4, "원인별 평균 정체길이 (km)", 8.3, 4.25, 4.5, 0.4, size=13, bold=True, color=DGRAY)
for i,(name,cnt,pct,dur,avg_l,col) in enumerate(causes):
    by = 4.75 + i * 0.4
    add_text(s4, name, 8.3, by, 1.7, 0.35, size=10, color=DGRAY)
    bar_h(s4, 10.1, by+0.04, 2.5, 0.26, avg_l, 5.0, color=col)
    add_text(s4, f"{avg_l}km", 12.65, by, 0.7, 0.35, size=10, bold=True, color=col)

# 인사이트
add_rect(s4, 0.35, 7.0, 12.6, 0.3, fill=RGBColor(0xFF, 0xF3, 0xE0), line=ORANGE, line_w=Pt(1))
add_text(s4,
    "행사는 평균 4.9km로 가장 넓은 정체 유발. 주말 정체는 건수 최소이나 지속시간 73분으로 최장 — 예측 대응 필요.",
    0.55, 7.02, 12.2, 0.28, size=10, color=DGRAY)

# ════════════════════════════════════════════════════════
# 슬라이드 5 — 시간대 · 요일별 분석
# ════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(s5, "시간대 · 요일별 분석", "정체 발생 시간 패턴 분석")
add_rect(s5, 0, 1.1, 13.33, 0.04, fill=BLUE)

# 시간대 데이터
time_bands = [
    ("심야\n00-06시", 32, 44.4, NAVY),
    ("출근\n06-10시", 12, 77.6, RED),
    ("오전\n10-14시", 34, 48.7, BLUE),
    ("오후\n14-18시", 16, 88.4, ORANGE),
    ("저녁\n18-22시", 14, 81.7, RGBColor(0x8E,0x44,0xAD)),
    ("야간\n22-24시", 12, 53.8, SKYB),
]

# 왼쪽 패널: 시간대
add_rect(s5, 0.35, 1.2, 8.3, 5.6, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_text(s5, "시간대별 발생 건수 및 평균 지속시간", 0.55, 1.25, 7, 0.4, size=13, bold=True, color=DGRAY)

for i,(band,cnt,dur,col) in enumerate(time_bands):
    bx = 0.55 + i * 1.35
    # 발생 건수 바 (세로형)
    bar_max_h = 3.2
    bar_h_val = bar_max_h * cnt / 34
    by_top = 5.0 - bar_h_val
    add_rect(s5, bx+0.1, 5.0-bar_max_h, 0.9, bar_max_h, fill=LGRAY)
    add_rect(s5, bx+0.1, by_top, 0.9, bar_h_val, fill=col)
    add_text(s5, f"{cnt}건", bx+0.05, by_top-0.38, 1.1, 0.35, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s5, band, bx, 5.05, 1.2, 0.55, size=9, color=DGRAY, align=PP_ALIGN.CENTER)
    add_text(s5, f"평균\n{dur}분", bx, 5.6, 1.2, 0.55, size=8, color=RGBColor(0x80,0x80,0x95), align=PP_ALIGN.CENTER)

add_text(s5, "← 건수", 0.38, 1.75, 0.8, 0.3, size=8, color=DGRAY)

# 오른쪽 패널: 요일
days_data = [
    ("월요일", 18, 55.8),
    ("화요일", 17, 54.6),
    ("수요일", 11, 82.8),
    ("목요일", 21, 54.0),
    ("금요일", 19, 54.4),
    ("토요일", 17, 67.3),
    ("일요일", 17, 62.1),
]
add_rect(s5, 8.9, 1.2, 4.1, 5.6, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_text(s5, "요일별 발생 건수", 9.1, 1.25, 3.6, 0.4, size=13, bold=True, color=DGRAY)
for i,(day,cnt,dur) in enumerate(days_data):
    by = 1.75 + i * 0.7
    col = RGBColor(0xE8,0x6A,0x2B) if day in ("토요일","일요일") else BLUE
    add_text(s5, day, 9.1, by+0.1, 1.1, 0.35, size=11, color=DGRAY)
    bar_h(s5, 10.3, by+0.1, 2.1, 0.35, cnt, 21, color=col)
    add_text(s5, f"{cnt}건", 12.45, by+0.1, 0.5, 0.35, size=11, bold=True, color=col)

# 인사이트
add_rect(s5, 0.35, 7.0, 12.6, 0.3, fill=RGBColor(0xE8,0xF5,0xE9), line=GREEN, line_w=Pt(1))
add_text(s5,
    "오전(10-14시) 34건으로 최다 발생. 오후(14-18시) 평균 88분으로 지속시간 최장. 수요일 82분으로 주중 최장 — 중반 집중 모니터링 필요.",
    0.55, 7.02, 12.2, 0.28, size=10, color=DGRAY)

# ════════════════════════════════════════════════════════
# 슬라이드 6 — 빈발 구간 & 최악의 정체
# ════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(s6, "주요 구간 분석", "빈발 정체 구간 TOP 8 및 최악의 정체 사례 TOP 5")
add_rect(s6, 0, 1.1, 13.33, 0.04, fill=BLUE)

# 왼쪽: 빈발 구간 TOP 8
add_rect(s6, 0.35, 1.2, 7.5, 5.9, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_rect(s6, 0.35, 1.2, 7.5, 0.45, fill=NAVY)
add_text(s6, "빈발 정체 구간 TOP 8 (발생 건수 기준)", 0.55, 1.24, 7, 0.38, size=12, bold=True, color=WHITE)

top_sections = [
    ("수원IC-안성IC",  "경부고속도로",  15, 53.3, 3.4),
    ("천안IC-대전IC",  "경부고속도로",  14, 41.1, 3.7),
    ("하남IC-청주JC",  "중부고속도로",  13, 53.8, 3.5),
    ("목감IC-당진IC",  "서해안고속도로",12, 65.6, 5.0),
    ("신갈JC-이천IC",  "영동고속도로",  11, 67.5, 4.2),
    ("안산JC-목감IC",  "서해안고속도로",11, 73.8, 3.4),
    ("여주IC-강릉JC",  "영동고속도로",  11, 50.8, 2.5),
    ("안성IC-천안IC",  "경부고속도로",  11, 65.4, 4.3),
]
colors6 = [ORANGE, BLUE, BLUE, GREEN, SKYB, GREEN, SKYB, BLUE]
for i,(sec,route,cnt,dur,avg_l) in enumerate(top_sections):
    by = 1.7 + i * 0.67
    bg = WHITE if i%2==0 else LGRAY
    add_rect(s6, 0.35, by, 7.5, 0.65, fill=bg)
    add_rect(s6, 0.35, by, 0.05, 0.65, fill=colors6[i])
    add_text(s6, f"#{i+1}", 0.45, by+0.12, 0.4, 0.4, size=11, bold=True, color=colors6[i])
    add_text(s6, sec,   0.9,  by+0.05, 2.0, 0.3, size=12, bold=True, color=DGRAY)
    add_text(s6, route, 0.9,  by+0.33, 2.2, 0.25, size=9, color=RGBColor(0x80,0x80,0x95))
    add_text(s6, f"{cnt}건", 3.1, by+0.12, 0.8, 0.4, size=14, bold=True, color=colors6[i])
    bar_h(s6, 3.95, by+0.2, 2.5, 0.28, cnt, 15, color=colors6[i])
    add_text(s6, f"평균 {dur}분", 6.55, by+0.12, 1.1, 0.3, size=10, color=DGRAY)

# 오른쪽: 최악의 정체 TOP 5
add_rect(s6, 8.1, 1.2, 5.0, 5.9, fill=WHITE, line=RGBColor(0xD0,0xD8,0xE8), line_w=Pt(1))
add_rect(s6, 8.1, 1.2, 5.0, 0.45, fill=RED)
add_text(s6, "최악의 정체 TOP 5 (지속시간 기준)", 8.3, 1.24, 4.6, 0.38, size=12, bold=True, color=WHITE)

worst = [
    ("04-09 18:45", "영동고속도로", "신갈JC-이천IC", "행사",    180, 12.9),
    ("04-08 08:00", "남해고속도로", "순천IC-광양IC", "행사",    168,  3.5),
    ("04-22 18:00", "경부고속도로", "수원IC-안성IC", "사고",    168, 14.9),
    ("04-02 17:00", "영동고속도로", "이천IC-여주IC", "주말 정체",165,  3.2),
    ("04-06 17:15", "서해안고속도로","목감IC-당진IC","기상악화", 156,  7.6),
]
for i,(dt,route,sec,cause,dur,km) in enumerate(worst):
    by = 1.75 + i * 1.05
    add_rect(s6, 8.1, by, 5.0, 1.0, fill=LGRAY if i%2==0 else WHITE)
    add_rect(s6, 8.1, by, 0.05, 1.0, fill=RED)
    badge_col = ORANGE if cause=="행사" else (RED if cause=="사고" else BLUE)
    add_rect(s6, 12.4, by+0.1, 0.65, 0.3, fill=badge_col)
    add_text(s6, cause, 12.4, by+0.1, 0.65, 0.3, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s6, f"{dur}분",  8.25, by+0.08, 1.0, 0.4, size=20, bold=True, color=RED)
    add_text(s6, dt,          8.25, by+0.52, 2.0, 0.3, size=9, color=RGBColor(0x80,0x80,0x95))
    add_text(s6, route,       10.0, by+0.08, 2.3, 0.28, size=10, bold=True, color=DGRAY)
    add_text(s6, sec,         10.0, by+0.35, 2.3, 0.28, size=10, color=DGRAY)
    add_text(s6, f"{km}km",   10.0, by+0.62, 2.3, 0.28, size=9, color=RGBColor(0x80,0x80,0x95))

# 인사이트
add_rect(s6, 0.35, 7.0, 12.6, 0.3, fill=RGBColor(0xFC, 0xE4, 0xEC), line=RED, line_w=Pt(1))
add_text(s6,
    "경부고속도로 수원IC-안성IC 구간이 15건으로 최다. 최악의 정체 5건 중 행사 원인이 2건 — 대형 행사 사전 교통 분산 대책 필수.",
    0.55, 7.02, 12.2, 0.28, size=10, color=DGRAY)

# ════════════════════════════════════════════════════════
# 슬라이드 7 — 종합 시사점 및 제언
# ════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(s7, "종합 시사점 및 제언", "분석 결과 기반 핵심 개선 방향")
add_rect(s7, 0, 1.1, 13.33, 0.04, fill=BLUE)

findings = [
    (BLUE,   "01", "경부고속도로 집중 관리",
     "전체의 40%(48건) 집중 발생\n수원IC-안성IC·천안IC-대전IC 구간 상시 모니터링 체계 강화",
     "수원IC-안성IC 15건 / 천안IC-대전IC 14건"),
    (ORANGE, "02", "행사·주말 사전 교통 분산",
     "행사 발생 시 평균 정체길이 4.9km — 원인별 최대\n주말 정체 평균 지속시간 73분으로 가장 오래 지속",
     "행사 원인 최악의 정체 5건 중 2건 차지"),
    (RED,    "03", "오후·저녁 시간대 집중 대응",
     "오후(14-18시) 평균 지속시간 88분 — 시간대별 최장\n저녁(18-22시)도 82분으로 퇴근 정체 장기화",
     "오전(10-14시) 34건으로 발생 빈도 최다"),
    (GREEN,  "04", "남해고속도로 발생 시 즉각 대응",
     "건수는 9건으로 최소이나 평균 지속 74분으로 최장\n발생 즉시 우회 안내 및 현장 투입 프로세스 필요",
     "평균 지속시간 전 노선 중 1위"),
]

for i,(col,num,title,body,stat) in enumerate(findings):
    row = i // 2
    c   = i %  2
    bx  = 0.35 + c * 6.5
    by  = 1.25 + row * 2.9

    add_rect(s7, bx, by, 6.2, 2.65, fill=WHITE, line=col, line_w=Pt(2))
    add_rect(s7, bx, by, 6.2, 0.08, fill=col)

    # 번호 배지
    add_rect(s7, bx+0.2, by+0.2, 0.55, 0.55, fill=col)
    add_text(s7, num, bx+0.2, by+0.2, 0.55, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s7, title, bx+0.9, by+0.25, 5.1, 0.45, size=14, bold=True, color=DGRAY)
    add_text(s7, body,  bx+0.2, by+0.82, 5.8, 1.0,  size=11, color=DGRAY)

    # 통계 박스
    add_rect(s7, bx+0.2, by+1.95, 5.8, 0.52, fill=RGBColor(0xF0,0xF4,0xF8))
    add_text(s7, f"DATA  {stat}", bx+0.35, by+2.05, 5.5, 0.32, size=10, color=col, bold=False)

# 하단 마무리
add_rect(s7, 0.35, 7.05, 12.6, 0.3, fill=NAVY)
add_text(s7,
    "본 분석은 2026년 4월 더미 데이터 기반이며, 실제 운영 데이터 적용 시 결과가 상이할 수 있습니다.",
    0.55, 7.07, 12.0, 0.25, size=9, color=RGBColor(0xA0,0xC0,0xE0), align=PP_ALIGN.CENTER)

# ── 저장 ────────────────────────────────────────────────
out = r"C:\Users\L2404161\Desktop\도로공사_바이브코딩\노선별_정체구간_분석보고서.pptx"
prs.save(out)
print(f"저장 완료: {out}")
